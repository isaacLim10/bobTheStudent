#####Imports#####
from requests import get
from bs4 import *
import re
from pptx import Presentation
from selenium import webdriver
import urllib.request

#####Globals#####
bulletLimit = 3


#####Functions#####
def createPowerpoint(filename, titleText, subtitleText, bullets, imageSlide):
    """
    This function creates a Powerpoint Presentation using pptx based on the params given.
    This code utilizes the OrderedDict class defined later
    :param filename: path of the file to write to(ie. article.pptx)
    :param titleText: title of the entire presentation(ie. Darth Maul)
    :param subtitleText: Subtitle of the entire presentation(ie. Generated with Python)
    :param bullets: A reference to an OrderedDict object with the bullet points and slide titles ie Overview : ['Content[0]', 'Content[1]'] etc.,
    :param imageSlide: The path of the image file to put onto the title slide
    """
    prs = Presentation() #Creates presentation object
    title_slide_layout = prs.slide_layouts[8] #Basically prs.slide_layouts is a list of all slide templates, with the title slide being last
    bullet_slide_layout = prs.slide_layouts[1] #And bullet slide is slide 2

    #Manipulating Title slide
    title_slide = prs.slides.add_slide(title_slide_layout) #Clones the title slide template and adds it to the actual presentation
    title = title_slide.shapes.title #Gets title placeholder
    subtitle = title_slide.placeholders[2] #Gets the subtitle placeholder
    image = title_slide.placeholders[1] #Gets the image

    title.text = titleText #Replaces title.text with what we want to be the titleText
    subtitle.text = subtitleText #Replaces subtitle.text with what we want
    image.insert_picture(imageSlide) #Inserts imageSlide picture into the image placeholder


    #Manipulating bullet slides
    for bulletTitle in bullets.keys: #Right so this iterates through the OrderedDict class and gets the bulletTitle each time
        slide = prs.slides.add_slide(bullet_slide_layout) #Adds this new slide to the presentation doc
        shapes = slide.shapes #All the placeholders

        title_shape = shapes.title #Gets title placeholder

        title_shape.text = bulletTitle #Replaces title placeholder with the current bulletTitle
        body_shape = shapes.placeholders[1].text_frame #Body block(ie. where all the bullets go)
        for bullet in bullets[bulletTitle]: #Every bullet within the bulletTitle
            tf = body_shape.add_paragraph() #Adds paragraph to the body
            tf.text = bullet #Changes the text on the paragraphs to the bullet

    prs.save(filename) #Saves the powerpoint to the filename path


def parseWikipedia(soupToParse):
    """
    Some regex code I stole off Stack Overflow(https://stackoverflow.com/questions/42324466/python-regular-expression-to-remove-all-square-brackets-and-their-contents)
    This basically takes a BeautifulSoup Element object and returns its text attribute, minus brackets and everything within
    :param soupToParse: Element object to parse
    :return: Its text attributes
    """
    pattern = r'\[.*?\]'
    return re.sub(pattern, '', soupToParse.text)

def splitLists(data, size):
    """
    More code I stole off stack(https://stackoverflow.com/questions/9671224/split-a-python-list-into-other-sublists-i-e-smaller-lists)
    Basically takes a list and splits it into sublists of size size
    :param data: Raw data to split
    :param size: Size of sublists
    :return: Split sublists
    """
    chunks = [data[x:x+size] for x in range(0, len(data), size)]
    return chunks

def getGoogleThumbnail(searchTerm, filepath):
    """
    Steals first thumbnail of searchTerm rom Google Images using selenium
    :param searchTerm: thing to search for
    :param filepath: path to save it to
    """
    driver = webdriver.Chrome("/Applications/PyCharm Edu.app/Contents/bin/Users/isaac_lims_macbook_air/PycharmProjects/feynman/chromedriver")
    driver.get("https://images.google.com/?gws_rd=ssl")
    searchBox = driver.find_element_by_css_selector("input#lst-ib.gsfi")
    searchBox.send_keys(searchTerm + "\n")
    imgThumbnail = driver.find_element_by_css_selector("img.rg_ic.rg_i")
    link = imgThumbnail.get_attribute("src")
    file = open(filepath, "wb")
    file.write(urllib.request.urlopen(link).read())
    file.close()




#####Classes#####
class Article(object):
    """
    This is a class to represent a Wikipedia Article.
    """
    def __init__(self, link):
        """
        :param link: The link of the Wikipedia Article
        """
        self.link = link
        doc = get(self.link)
        self.doc = doc.text
        self.headers = []
        self.paras = []
        self.soup = BeautifulSoup(self.doc, "lxml")

    def getHeaders(self):
        """
        :return: All h2 tags in the Article as bs4 elements
        """
        headers = self.soup.findAll('h2')
        for header in headers:
            self.headers.append(parseWikipedia(header))
        return self.headers

    def getParas(self):
        """
        :return: All p tags in the Article as bs4 elements
        """
        paras = self.soup.findAll('p')
        for para in paras:
            self.paras.append(parseWikipedia(para))
        return self.paras

class OrderedDict(object):
    """
    This is a special data structure I wanted to make.
    Basically I wanted a dictionary but with an order so when you iterate it goes in the order you put the keys in.
    It's not perfect but it's decent for my needs.
    Also this has SO MANY MAGIC METHODS.
    """
    def __init__(self):
        self.keys = []
        self.items = []

    def __getitem__(self, item):
        return self.items[self.keys.index(item)]

    def __setitem__(self, key, value):
        """
        This is a magic method that utilizes Python's default indexing system. ie square brackets
        Take example dictName["Meaning Of Life"] = 42
        :param key: The key here is "Meaning Of Life"
        :param value: The value is 42
        """
        """
        What I'm doing here is trying to change get the position of the key in the self.keys list using the index method.
        This will raise an error if the key stated is not already in the self.keys list. To deal with this, I have
        created an except route where the item and key is appended at the same time so that they have the same index(assuming
        they start off with the same index which should be true.
        """
        try:
            keyPos = self.keys.index(key) #Getting position of key in self.keys then
            self.items[keyPos] = value #Setting the item at that position in the self.items list to the value so they have the same index
        except ValueError: #Escape route for if appending a new thing
            self.items.append(value)
            self.keys.append(key)

    def __iter__(self): #Iterator function that iterates through a list
        self.n = -1 #Self.n starts at -1 because it will be immediately incremented
        return self #I couldn't figure out why or how this work but it seems to so...

    def __next__(self): #Used for the iterator function to progress
        if self.n <= len(self) -2: #Checks if self.n is lower than the number items in this object(-2 is because it starts at -1 instead of 0)
            self.n += 1 #Increments self.n
            return self.items[self.n]
        else:
            raise StopIteration #Stops the iteration

    def __len__(self):
        """
        Magic method for when you do len(object)
        :return: length of the "dictionary"
        """
        return len(self.keys) #Returns length of self.keys

class OrderedDictError(Exception):pass #Custom exception


if __name__ == "__main__": #If this is running straight from this file
    searchTerm = input("What to search for:")
    link = "https://en.wikipedia.org/wiki/" + searchTerm.replace(" ", "_")
    getGoogleThumbnail(searchTerm, "article.png")
    article = Article(link)
    article.getParas()
    article.getHeaders()
    splitParas = (splitLists(article.paras, bulletLimit)) #Splits the paragraphs into sublists of 5 each
    slides = OrderedDict() #Creates orderedDict object
    slides["Overveew"] = article.headers
    for i in range(len(splitParas)):
        slides["Body(" + str(i + 1) + ")"] = splitParas[i]
    slides["Bibiography"] = link.replace("\n", "") #Terrible Bibliography
    createPowerpoint("article.pptx", searchTerm, "Generated with Python", slides, "article.png") #daveed.png is a placeholder image(Daveed Diggs is amazing)
